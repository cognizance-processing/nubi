"""Tests for app/embedding/render_pptx.py — T4 PPTX renderer.

Coverage
--------
1. render_board_pptx — small board (chart + kpi + text) with pre-rendered
   SVGs produces a valid .pptx: correct slide count, at least one picture
   shape per data widget.

2. render_board_pptx — title_slide config inserts an extra slide at position 0.

3. render_board_pptx — widget excluded via include_in_report=False does not
   appear as a slide.

4. render_board_pptx — caption is added to the slide shapes (text box with
   caption text).

5. render_board_pptx — header/footer strings appear on each widget slide.

6. render_board_pptx — page_size='16:9' sets the correct slide dimensions.

7. render_board_pptx — widget with svg=None (WebGL) still produces a slide
   (placeholder text box, no picture shape).

All tests are offline — no live DB, no Node.js. python-pptx + cairosvg must
be available (tests are skipped otherwise, matching the optional-dep convention).
"""

from __future__ import annotations

import io
from typing import Any

import pytest

# Skip entire module when python-pptx is not installed.
pptx = pytest.importorskip("pptx", reason="python-pptx not installed — skipping PPTX render tests")

# cairosvg requires the libcairo system library; skip when it cannot load.
try:
    import cairosvg as _cairosvg  # noqa: F401
    _CAIROSVG_AVAILABLE = True
except (ImportError, OSError):
    _CAIROSVG_AVAILABLE = False

pytestmark = pytest.mark.skipif(
    not _CAIROSVG_AVAILABLE,
    reason="cairosvg / libcairo not available — skipping PPTX render tests",
)

from pptx import Presentation

from app.dashboards.spec import (
    BoardExportConfig,
    DashboardSpec,
    SurfaceGridEntry,
    Surfaces,
    TitleSlideConfig,
    Widget,
    WidgetExportHints,
    WidgetPos,
)
from app.embedding.render_pptx import render_board_pptx


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

_SIMPLE_SVG = (
    '<svg xmlns="http://www.w3.org/2000/svg" width="400" height="200">'
    '<rect width="400" height="200" fill="#e0e0e0"/>'
    '<text x="20" y="100" font-size="24">Widget</text>'
    "</svg>"
)


def _small_spec(
    *,
    title_slide: TitleSlideConfig | None = None,
    header: str | None = None,
    footer: str | None = None,
    page_size: str = "A4",
    widget_hints: dict[str, Any] | None = None,
) -> DashboardSpec:
    """Build a small 3-widget spec (chart + kpi + text) with export config."""
    wh: dict[str, WidgetExportHints] = {}
    if widget_hints:
        for wid, hints_dict in widget_hints.items():
            wh[wid] = WidgetExportHints(**hints_dict)

    return DashboardSpec(
        title="Test Board",
        layout={"cols": 12, "row_height": 60},
        widgets=[
            Widget(
                id="chart1",
                type="chart",
                query_id="q1",
                chart_type="bar",
                encoding={"x": "region", "y": "revenue"},
                pos=WidgetPos(x=1, y=1, w=6, h=3),
            ),
            Widget(
                id="kpi1",
                type="kpi",
                query_id="q1",
                encoding={"value": "revenue"},
                props={"label": "Total Revenue"},
                pos=WidgetPos(x=7, y=1, w=3, h=2),
            ),
            Widget(
                id="txt1",
                type="text",
                content="## Hello\nThis is a text widget.",
                pos=WidgetPos(x=10, y=1, w=3, h=2),
            ),
        ],
        export=BoardExportConfig(
            title_slide=title_slide,
            header=header,
            footer=footer,
            page_size=page_size,  # type: ignore[arg-type]
            widget_hints=wh,
        ),
    )


def _parse_pptx(raw_bytes: bytes) -> Presentation:
    """Load raw pptx bytes into a Presentation object for assertions."""
    return Presentation(io.BytesIO(raw_bytes))


def _count_picture_shapes(slide: Any) -> int:
    """Return the number of picture shapes on *slide*."""
    from pptx.shapes.base import BaseShape
    count = 0
    for shape in slide.shapes:
        # A picture shape has shape_type == MSO_SHAPE_TYPE.PICTURE (13)
        if shape.shape_type == 13:
            count += 1
    return count


def _text_in_slide(slide: Any) -> list[str]:
    """Return all non-empty text strings from text boxes + placeholders on *slide*."""
    texts: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    return texts


# ---------------------------------------------------------------------------
# 1. Basic render — 3 widgets → 3 slides, each with a picture shape
# ---------------------------------------------------------------------------


def test_render_board_pptx_basic_slide_count():
    """3-widget board produces a 3-slide PPTX with a picture on each slide."""
    spec = _small_spec()
    svgs = {
        "chart1": _SIMPLE_SVG,
        "kpi1": _SIMPLE_SVG,
        "txt1": _SIMPLE_SVG,
    }

    raw = render_board_pptx(spec, svgs)

    assert isinstance(raw, bytes)
    assert len(raw) > 0

    prs = _parse_pptx(raw)
    assert len(prs.slides) == 3, f"Expected 3 slides, got {len(prs.slides)}"

    for slide in prs.slides:
        pics = _count_picture_shapes(slide)
        assert pics >= 1, f"Expected at least 1 picture per slide, got {pics}"


# ---------------------------------------------------------------------------
# 2. Title slide inserts an extra slide at position 0
# ---------------------------------------------------------------------------


def test_render_board_pptx_title_slide():
    """title_slide config adds an extra title slide as the first slide."""
    spec = _small_spec(
        title_slide=TitleSlideConfig(heading="Q1 Report", subheading="Confidential"),
    )
    svgs = {"chart1": _SIMPLE_SVG, "kpi1": _SIMPLE_SVG, "txt1": _SIMPLE_SVG}

    raw = render_board_pptx(spec, svgs)
    prs = _parse_pptx(raw)

    # 1 title slide + 3 widget slides = 4
    assert len(prs.slides) == 4, f"Expected 4 slides (1 title + 3 widgets), got {len(prs.slides)}"

    # Title slide text should contain the heading.
    title_texts = _text_in_slide(prs.slides[0])
    assert any("Q1 Report" in t for t in title_texts), (
        f"Expected heading 'Q1 Report' in title slide, got: {title_texts}"
    )


# ---------------------------------------------------------------------------
# 3. include_in_report=False excludes a widget
# ---------------------------------------------------------------------------


def test_render_board_pptx_exclude_widget():
    """A widget with include_in_report=False is not rendered as a slide."""
    spec = _small_spec(
        widget_hints={
            "kpi1": {"include_in_report": False},
        }
    )
    svgs = {"chart1": _SIMPLE_SVG, "kpi1": _SIMPLE_SVG, "txt1": _SIMPLE_SVG}

    raw = render_board_pptx(spec, svgs)
    prs = _parse_pptx(raw)

    # Only chart1 + txt1 should produce slides (kpi1 excluded).
    assert len(prs.slides) == 2, (
        f"Expected 2 slides (kpi1 excluded), got {len(prs.slides)}"
    )


# ---------------------------------------------------------------------------
# 4. Caption appears as text on the slide
# ---------------------------------------------------------------------------


def test_render_board_pptx_caption():
    """A widget hint with caption produces a text box on the slide."""
    spec = _small_spec(
        widget_hints={
            "chart1": {"caption": "Figure 1: Revenue by Region"},
        }
    )
    svgs = {"chart1": _SIMPLE_SVG, "kpi1": _SIMPLE_SVG, "txt1": _SIMPLE_SVG}

    raw = render_board_pptx(spec, svgs)
    prs = _parse_pptx(raw)

    # chart1 is the first widget → first slide.
    first_slide_texts = _text_in_slide(prs.slides[0])
    assert any("Figure 1" in t for t in first_slide_texts), (
        f"Expected caption 'Figure 1: Revenue by Region' on first slide, "
        f"got: {first_slide_texts}"
    )


# ---------------------------------------------------------------------------
# 5. Header and footer text appears on each widget slide
# ---------------------------------------------------------------------------


def test_render_board_pptx_header_footer():
    """Header/footer strings appear as text boxes on each widget slide."""
    spec = _small_spec(header="Nubi Analytics", footer="Confidential")
    svgs = {"chart1": _SIMPLE_SVG, "kpi1": _SIMPLE_SVG, "txt1": _SIMPLE_SVG}

    raw = render_board_pptx(spec, svgs)
    prs = _parse_pptx(raw)

    for i, slide in enumerate(prs.slides):
        texts = _text_in_slide(slide)
        assert any("Nubi Analytics" in t for t in texts), (
            f"Slide {i}: expected header 'Nubi Analytics', got: {texts}"
        )
        assert any("Confidential" in t for t in texts), (
            f"Slide {i}: expected footer 'Confidential', got: {texts}"
        )


# ---------------------------------------------------------------------------
# 6. page_size='16:9' sets the correct slide dimensions
# ---------------------------------------------------------------------------


def test_render_board_pptx_page_size_16_9():
    """page_size='16:9' sets slide dimensions to the widescreen preset."""
    from pptx.util import Cm

    spec = _small_spec(page_size="16:9")
    svgs = {"chart1": _SIMPLE_SVG, "kpi1": _SIMPLE_SVG, "txt1": _SIMPLE_SVG}

    raw = render_board_pptx(spec, svgs)
    prs = _parse_pptx(raw)

    expected_w_emu = Cm(33.87)
    expected_h_emu = Cm(19.05)

    # Allow a small rounding tolerance (1 EMU).
    assert abs(prs.slide_width - expected_w_emu) <= 1, (
        f"Expected slide width ~{expected_w_emu} EMU, got {prs.slide_width}"
    )
    assert abs(prs.slide_height - expected_h_emu) <= 1, (
        f"Expected slide height ~{expected_h_emu} EMU, got {prs.slide_height}"
    )


# ---------------------------------------------------------------------------
# 7. WebGL widget (svg=None) produces a slide with a placeholder text box
# ---------------------------------------------------------------------------


def test_render_board_pptx_webgl_placeholder():
    """A widget with svg=None produces a slide without a picture shape."""
    spec = _small_spec()
    svgs: dict[str, str | None] = {
        "chart1": None,   # WebGL widget — no SVG
        "kpi1": _SIMPLE_SVG,
        "txt1": _SIMPLE_SVG,
    }

    raw = render_board_pptx(spec, svgs)
    prs = _parse_pptx(raw)

    assert len(prs.slides) == 3, f"Expected 3 slides, got {len(prs.slides)}"

    # First slide (chart1) has NO picture shapes.
    chart_pics = _count_picture_shapes(prs.slides[0])
    assert chart_pics == 0, (
        f"Expected 0 picture shapes for WebGL placeholder slide, got {chart_pics}"
    )

    # First slide has a text box (the placeholder message).
    chart_texts = _text_in_slide(prs.slides[0])
    assert any("chart1" in t.lower() or "webgl" in t.lower() for t in chart_texts), (
        f"Expected WebGL placeholder text on first slide, got: {chart_texts}"
    )

    # Subsequent slides (kpi1, txt1) have picture shapes.
    for i in [1, 2]:
        pics = _count_picture_shapes(prs.slides[i])
        assert pics >= 1, f"Expected picture on slide {i + 1}, got {pics}"
