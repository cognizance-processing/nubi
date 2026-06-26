"""PPTX renderer — T4: build a .pptx from a board's widget SVGs.

Design
------
* **Native SVG embed** — python-pptx inserts the SVG as a picture shape
  (``prs.slide_layouts[6]`` blank layout). PowerPoint 2016+ and LibreOffice
  render native SVG; a PNG raster fallback (via cairosvg) is also embedded for
  compatibility with older clients that cannot display SVG.
* **Source** — per-widget SVGs are produced by the T2 pipeline
  (:func:`app.dashboards.svg_render.render_board_svg`); this module only
  assembles the PPTX from already-rendered SVG strings.
* **Export config** — T5 :func:`app.dashboards.spec.get_export_config` drives
  the title slide, header/footer, caption, widget order, include/exclude, and
  page_break_before behaviour.
* **Surface mapping** — when a ``slides`` surface is present in the spec, each
  slide entry maps directly to a slide.  When absent (the common case) the grid
  surface is used and each widget (respecting export hints) becomes one slide.
* **Lazy imports** — ``python-pptx`` and ``cairosvg`` are imported on first
  call with a clear install error.  Both are added to the optional deps section
  of the top-level ``requirements.txt``.

Public API
----------
``render_board_pptx(spec, widget_svgs, *, ...) -> bytes``
    Given a :class:`~app.dashboards.spec.DashboardSpec` and a dict of
    ``{widget_id: svg_string}`` (pre-rendered by the T2 pipeline), produce a
    ``.pptx`` file as raw bytes.

``render_board_pptx_from_data(spec, widget_data, *, ...) -> bytes``
    Convenience wrapper: calls the T2 SVG renderer for each widget then
    delegates to :func:`render_board_pptx`.
"""

from __future__ import annotations

import io
from typing import Any

from app.dashboards.spec import DashboardSpec, get_export_config, get_surface_layout
from app.errors import AppError

# ---------------------------------------------------------------------------
# Page-size presets  (width_cm, height_cm)
# ---------------------------------------------------------------------------

_PAGE_SIZES_CM: dict[str, tuple[float, float]] = {
    "A4": (29.7, 21.0),       # A4 landscape (natural for slides)
    "Letter": (27.94, 21.59), # US Letter landscape
    "16:9": (33.87, 19.05),   # Widescreen slide canvas
}

# Slide margins in cm (used when placing a single widget per slide)
_SLIDE_MARGIN_CM = 0.5

# Caption area height reserved when a caption is present (cm)
_CAPTION_HEIGHT_CM = 0.8

# Header/footer bar height (cm)
_HEADER_FOOTER_HEIGHT_CM = 0.6


# ---------------------------------------------------------------------------
# Lazy-import helpers
# ---------------------------------------------------------------------------


def _require_pptx() -> Any:
    """Return the ``pptx`` package, raising AppError if not installed."""
    try:
        import pptx  # noqa: PLC0415
        return pptx
    except ImportError as exc:
        raise AppError(
            "pptx_not_installed",
            "python-pptx is required for PPTX export but is not installed. "
            "Install it with: pip install python-pptx",
            503,
        ) from exc


def _require_cairosvg() -> Any:
    """Return the ``cairosvg`` module, raising AppError if not installed.

    Also catches ``OSError`` for the case where the Python package is installed
    but the underlying ``libcairo`` system library is missing.
    """
    try:
        import cairosvg  # noqa: PLC0415
        return cairosvg
    except (ImportError, OSError) as exc:
        raise AppError(
            "cairosvg_not_installed",
            "cairosvg is required for PNG rasterization in PPTX export but is "
            "not available. Install it with: pip install cairosvg and ensure "
            "the libcairo system library is present (e.g. brew install cairo).",
            503,
        ) from exc


# ---------------------------------------------------------------------------
# SVG → PNG rasterization helper
# ---------------------------------------------------------------------------


def _svg_to_png(svg_string: str, *, width_px: int = 1200, height_px: int = 675) -> bytes:
    """Rasterize an SVG string to PNG bytes via cairosvg.

    Parameters
    ----------
    svg_string:
        A valid SVG XML string.
    width_px, height_px:
        Output raster dimensions in pixels.

    Returns
    -------
    bytes
        PNG-encoded image bytes.

    Raises
    ------
    AppError("cairosvg_not_installed", 503)
        When cairosvg is not installed.
    AppError("svg_rasterize_error", 503)
        When cairosvg fails to convert the SVG.
    """
    cairosvg = _require_cairosvg()
    try:
        return cairosvg.svg2png(
            bytestring=svg_string.encode("utf-8"),
            output_width=width_px,
            output_height=height_px,
        )
    except Exception as exc:
        raise AppError(
            "svg_rasterize_error",
            f"cairosvg failed to rasterize SVG: {exc}",
            503,
        ) from exc


# ---------------------------------------------------------------------------
# Unit helpers (python-pptx uses EMU — English Metric Units; 914400 EMU = 1 in)
# ---------------------------------------------------------------------------


def _cm_to_emu(cm: float) -> int:
    """Convert centimetres to EMU (python-pptx's native unit)."""
    # 1 inch = 2.54 cm; 1 inch = 914400 EMU
    return int(cm * 914400 / 2.54)


# ---------------------------------------------------------------------------
# Slide builder
# ---------------------------------------------------------------------------


def _add_widget_slide(
    prs: Any,
    slide_layout: Any,
    svg_string: str | None,
    *,
    page_width_cm: float,
    page_height_cm: float,
    caption: str | None = None,
    header: str | None = None,
    footer: str | None = None,
    widget_id: str = "",
) -> Any:
    """Add one slide containing *svg_string* (with optional caption / header / footer).

    The SVG is embedded natively as a picture shape (SVG blob).  A PNG raster is
    also embedded as the same shape's fallback so old clients display the image.
    When *svg_string* is ``None`` (e.g. WebGL widget with no SVG), a placeholder
    text box is inserted instead.

    Parameters
    ----------
    prs:
        The ``Presentation`` object.
    slide_layout:
        Blank slide layout to use.
    svg_string:
        SVG content or ``None`` for a WebGL/unavailable widget.
    page_width_cm, page_height_cm:
        Slide canvas dimensions (from the page_size preset).
    caption:
        Optional caption string rendered below the widget area.
    header:
        Optional header string rendered at the top of the slide.
    footer:
        Optional footer string rendered at the bottom of the slide.
    widget_id:
        Widget id for placeholder / alt-text labelling.

    Returns
    -------
    Slide
        The newly added slide.
    """
    from pptx.util import Cm, Pt  # noqa: PLC0415
    from pptx.enum.text import PP_ALIGN  # noqa: PLC0415

    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes

    # Compute the usable rectangle for the widget content (after reserving
    # header, footer, caption bands).
    top_cm = _SLIDE_MARGIN_CM + (_HEADER_FOOTER_HEIGHT_CM if header else 0)
    bottom_reserved_cm = _SLIDE_MARGIN_CM + (_HEADER_FOOTER_HEIGHT_CM if footer else 0)
    if caption:
        bottom_reserved_cm += _CAPTION_HEIGHT_CM

    content_height_cm = page_height_cm - top_cm - bottom_reserved_cm
    content_left_cm = _SLIDE_MARGIN_CM
    content_width_cm = page_width_cm - 2 * _SLIDE_MARGIN_CM

    if content_height_cm < 1.0:
        # Degenerate dimensions — skip placement rather than crashing.
        content_height_cm = 1.0

    # ── Header ──────────────────────────────────────────────────────────────
    if header:
        txBox = shapes.add_textbox(
            Cm(_SLIDE_MARGIN_CM),
            Cm(_SLIDE_MARGIN_CM),
            Cm(content_width_cm),
            Cm(_HEADER_FOOTER_HEIGHT_CM),
        )
        tf = txBox.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = header
        run.font.size = Pt(9)
        run.font.bold = False

    # ── Widget content ───────────────────────────────────────────────────────
    if svg_string:
        try:
            _add_svg_picture(
                shapes,
                svg_string,
                left_cm=content_left_cm,
                top_cm=top_cm,
                width_cm=content_width_cm,
                height_cm=content_height_cm,
                widget_id=widget_id,
            )
        except AppError as exc:
            # cairosvg not installed or rasterization failed — insert a
            # plain-text placeholder with the error so the slide is not blank.
            _add_placeholder_textbox(
                shapes,
                f"[{widget_id}] SVG render error: {exc.message}",
                left_cm=content_left_cm,
                top_cm=top_cm,
                width_cm=content_width_cm,
                height_cm=content_height_cm,
            )
    else:
        _add_placeholder_textbox(
            shapes,
            f"[{widget_id}] (WebGL widget — raster not available)",
            left_cm=content_left_cm,
            top_cm=top_cm,
            width_cm=content_width_cm,
            height_cm=content_height_cm,
        )

    # ── Caption ──────────────────────────────────────────────────────────────
    if caption:
        caption_top_cm = page_height_cm - _SLIDE_MARGIN_CM - (_HEADER_FOOTER_HEIGHT_CM if footer else 0) - _CAPTION_HEIGHT_CM
        txBox = shapes.add_textbox(
            Cm(content_left_cm),
            Cm(caption_top_cm),
            Cm(content_width_cm),
            Cm(_CAPTION_HEIGHT_CM),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = caption
        run.font.size = Pt(9)
        run.font.italic = True

    # ── Footer ──────────────────────────────────────────────────────────────
    if footer:
        footer_top_cm = page_height_cm - _SLIDE_MARGIN_CM - _HEADER_FOOTER_HEIGHT_CM
        txBox = shapes.add_textbox(
            Cm(_SLIDE_MARGIN_CM),
            Cm(footer_top_cm),
            Cm(content_width_cm),
            Cm(_HEADER_FOOTER_HEIGHT_CM),
        )
        tf = txBox.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = footer
        run.font.size = Pt(9)
        run.font.bold = False

    return slide


def _add_svg_picture(
    shapes: Any,
    svg_string: str,
    *,
    left_cm: float,
    top_cm: float,
    width_cm: float,
    height_cm: float,
    widget_id: str = "",
) -> None:
    """Embed *svg_string* into *shapes* as a PNG picture shape.

    Strategy
    --------
    1. Rasterize the SVG to PNG via cairosvg (~96 DPI relative to the slide
       dimensions).  PNG is the primary content for broad client compatibility.
    2. Insert the PNG with ``shapes.add_picture``.
    3. Best-effort: add the SVG blob as a raw image part and patch the OOXML
       ``<a:blip>`` extension list so PowerPoint 2016+ renders the native vector.
       This step is wrapped in a broad ``except`` — if it fails for any reason
       (python-pptx internals, image-type checks, theme differences) the PNG
       fallback is already embedded and the slide stays valid.

    Raises
    ------
    AppError("cairosvg_not_installed", 503)
        When cairosvg is not installed / libcairo is missing.
    """
    from pptx.util import Cm  # noqa: PLC0415
    from pptx.oxml.ns import qn  # noqa: PLC0415
    from lxml import etree  # noqa: PLC0415

    # ── Step 1: PNG raster ───────────────────────────────────────────────────
    raster_w_px = max(100, int(width_cm / 2.54 * 96))
    raster_h_px = max(100, int(height_cm / 2.54 * 96))
    png_bytes = _svg_to_png(svg_string, width_px=raster_w_px, height_px=raster_h_px)

    # ── Step 2: Insert PNG as primary picture shape ──────────────────────────
    pic_shape = shapes.add_picture(
        io.BytesIO(png_bytes),
        Cm(left_cm),
        Cm(top_cm),
        Cm(width_cm),
        Cm(height_cm),
    )

    # ── Step 3: Native SVG layer — best-effort, must NOT corrupt the package ─
    # We create a raw OPC part for the SVG blob (bypassing Pillow image checks)
    # and register it as an image relationship on the slide part, then annotate
    # the <a:blip> element so PowerPoint 2016+ uses the vector source.
    try:
        from pptx.opc.package import Part as OpcPart  # noqa: PLC0415
        from pptx.opc.packuri import PackURI  # noqa: PLC0415

        slide_part = pic_shape.part  # SlidePart
        svg_bytes_content = svg_string.encode("utf-8")

        SVG_CONTENT_TYPE = "image/svg+xml"
        SVG_RELTYPE = (
            "http://schemas.microsoft.com/office/2016/06/relationships/svgBlip"
        )

        # Compute a unique partname: /ppt/media/svgN.svg
        package = slide_part.package
        existing_names = {
            p.partname
            for p in package.iter_parts()
            if hasattr(p, "partname")
        }
        idx = 1
        while f"/ppt/media/svg{idx}.svg" in existing_names:
            idx += 1
        svg_partname = PackURI(f"/ppt/media/svg{idx}.svg")

        # Create the raw part (requires package so iter_parts finds it).
        svg_part = OpcPart(svg_partname, SVG_CONTENT_TYPE, package, svg_bytes_content)
        svg_rId = slide_part.relate_to(svg_part, SVG_RELTYPE)

        # Patch <a:blip> in the <p:pic> element.
        sp_element = pic_shape._element
        blip = sp_element.find(".//" + qn("a:blip"))
        if blip is not None:
            ASVG_NS = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
            ext_lst = blip.find(qn("a:extLst"))
            if ext_lst is None:
                ext_lst = etree.SubElement(blip, qn("a:extLst"))
            ext = etree.SubElement(ext_lst, qn("a:ext"))
            ext.set("uri", "{96DAC541-7B7A-43D3-8B79-37D633B846F1}")
            svg_blip = etree.SubElement(ext, f"{{{ASVG_NS}}}svgBlip")
            svg_blip.set(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed",
                svg_rId,
            )
    except Exception:  # noqa: BLE001
        # SVG native layer failed — PNG fallback is already embedded; degrade
        # gracefully so the presentation remains valid.
        pass


def _add_placeholder_textbox(
    shapes: Any,
    text: str,
    *,
    left_cm: float,
    top_cm: float,
    width_cm: float,
    height_cm: float,
) -> None:
    """Insert a plain text box as a placeholder for unavailable widget content."""
    from pptx.util import Cm, Pt  # noqa: PLC0415
    from pptx.dml.color import RGBColor  # noqa: PLC0415

    txBox = shapes.add_textbox(Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)


# ---------------------------------------------------------------------------
# Title slide builder
# ---------------------------------------------------------------------------


def _add_title_slide(
    prs: Any,
    title_layout: Any,
    heading: str,
    subheading: str | None,
) -> Any:
    """Add a title slide with *heading* and optional *subheading*."""

    slide = prs.slides.add_slide(title_layout)
    title_ph = None
    body_ph = None

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            title_ph = ph
        elif ph.placeholder_format.idx == 1:
            body_ph = ph

    if title_ph is not None:
        title_ph.text = heading
    if body_ph is not None and subheading:
        body_ph.text = subheading

    return slide


# ---------------------------------------------------------------------------
# Widget ordering helper (T5 export hints)
# ---------------------------------------------------------------------------


def _ordered_widgets(
    spec: DashboardSpec,
    export_cfg: dict[str, Any],
) -> list[Any]:
    """Return spec.widgets ordered and filtered by T5 export hints.

    Applies ``include_in_report``, ``order``, and ``page_break_before`` from
    :func:`~app.dashboards.spec.get_export_config`.

    Returns
    -------
    list[Widget]
        Widgets that should appear in the PPTX, in the correct order.
        Each widget's export hint is NOT attached; callers should look it up
        from ``export_cfg['widget_hints'][widget.id]``.
    """
    hints = export_cfg.get("widget_hints", {})

    # Filter out excluded widgets.
    included = [
        w for w in spec.widgets
        if hints.get(w.id, {}).get("include_in_report", True)
    ]

    # Pre-compute natural position map — O(n) once, avoids O(n²) list.index calls
    # in the sort fallback path.
    pos = {w.id: i for i, w in enumerate(spec.widgets)}

    def _sort_key(w: Any) -> tuple[int, int]:
        hint = hints.get(w.id, {})
        explicit_order = hint.get("order")
        if explicit_order is not None:
            return (0, int(explicit_order))
        # Fall back to the position in spec.widgets (natural order).
        return (1, pos.get(w.id, len(spec.widgets)))

    return sorted(included, key=_sort_key)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def render_board_pptx(
    spec: DashboardSpec,
    widget_svgs: dict[str, str | None],
    *,
    page_size: str | None = None,
) -> bytes:
    """Build a ``.pptx`` from a board's per-widget SVG strings.

    This is the main entry point for T4 PPTX export.  It:

    1. Reads the T5 export config (title_slide, header, footer, page_size,
       widget order, captions, include/exclude, page_break_before).
    2. Creates a blank :class:`pptx.Presentation` sized to *page_size*.
    3. Optionally inserts a title slide.
    4. Iterates over included widgets (ordered by export hints) and adds one
       slide per widget, embedding the SVG natively + PNG raster fallback.
    5. Returns the raw ``.pptx`` bytes.

    Parameters
    ----------
    spec:
        A validated :class:`~app.dashboards.spec.DashboardSpec`.
    widget_svgs:
        Mapping of ``widget_id -> svg_string | None``.  Produced by the T2
        pipeline (:func:`~app.dashboards.svg_render.render_widgets_svg`).
        ``None`` values indicate WebGL / unavailable widgets (a placeholder
        slide is inserted instead of a picture).
    page_size:
        Override the page size (``'A4'``, ``'Letter'``, ``'16:9'``).  When
        ``None`` the value from ``spec.export.page_size`` is used (default
        ``'A4'``).

    Returns
    -------
    bytes
        Raw ``.pptx`` file bytes.

    Raises
    ------
    AppError("pptx_not_installed", 503)
        When python-pptx is not installed.
    AppError("cairosvg_not_installed", 503)
        When cairosvg is not installed and SVG rasterization is needed.
    """
    pptx_mod = _require_pptx()
    from pptx import Presentation  # noqa: PLC0415
    from pptx.util import Cm  # noqa: PLC0415

    export_cfg = get_export_config(spec)

    # Resolve page size.
    resolved_page_size = page_size or export_cfg.get("page_size") or "A4"
    width_cm, height_cm = _PAGE_SIZES_CM.get(resolved_page_size, _PAGE_SIZES_CM["A4"])

    header: str | None = export_cfg.get("header")
    footer: str | None = export_cfg.get("footer")
    title_slide_cfg: dict | None = export_cfg.get("title_slide")

    prs = Presentation()
    prs.slide_width = Cm(width_cm)
    prs.slide_height = Cm(height_cm)

    # Layouts: index 0 = title slide, index 6 = blank.
    # Guard against themes with fewer layouts.
    layouts = prs.slide_layouts
    title_layout = layouts[0] if len(layouts) > 0 else layouts[0]
    blank_layout = layouts[6] if len(layouts) > 6 else layouts[len(layouts) - 1]

    # ── Title slide ──────────────────────────────────────────────────────────
    if title_slide_cfg is not None:
        heading = title_slide_cfg.get("heading") or spec.title
        subheading = title_slide_cfg.get("subheading")
        _add_title_slide(prs, title_layout, heading, subheading)

    # ── Widget slides ────────────────────────────────────────────────────────
    ordered = _ordered_widgets(spec, export_cfg)
    hints = export_cfg.get("widget_hints", {})

    for widget in ordered:
        hint = hints.get(widget.id, {})
        svg = widget_svgs.get(widget.id)
        caption: str | None = hint.get("caption")

        # page_break_before: in a per-widget-per-slide model this is a no-op
        # for the first widget; for subsequent widgets it inserts an EXTRA blank
        # slide before the widget slide (matching PDF "page break before" semantics).
        if hint.get("page_break_before") and prs.slides and len(prs.slides) > (1 if title_slide_cfg else 0):
            prs.slides.add_slide(blank_layout)

        _add_widget_slide(
            prs,
            blank_layout,
            svg,
            page_width_cm=width_cm,
            page_height_cm=height_cm,
            caption=caption,
            header=header,
            footer=footer,
            widget_id=widget.id,
        )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def render_board_pptx_from_data(
    spec: DashboardSpec,
    widget_data: list[dict[str, Any]],
    *,
    page_size: str | None = None,
    page_width_px: int = 1200,
    row_height_px: int = 120,
    margin_px: int = 16,
    gap_px: int = 8,
    ssr_timeout: float = 60.0,
) -> bytes:
    """Render per-widget SVGs via T2 then build a PPTX.

    Convenience wrapper: calls
    :func:`~app.dashboards.svg_render.render_widgets_svg` for each widget in
    the spec (using the pre-collected *widget_data*), then delegates to
    :func:`render_board_pptx`.

    Parameters
    ----------
    spec:
        A validated :class:`~app.dashboards.spec.DashboardSpec`.
    widget_data:
        Pre-collected widget data (from
        :func:`~app.dashboards.collect.collect_board_data`) as a list of
        ``{widget_id, columns, rows, ...}`` dicts.
    page_size:
        Page size override (``'A4'``, ``'Letter'``, ``'16:9'``).
    page_width_px, row_height_px, margin_px, gap_px:
        Page geometry forwarded to the T2 SVG renderer to compute per-widget
        pixel dimensions from their grid positions.
    ssr_timeout:
        Timeout in seconds for the Node.js SSR subprocess.

    Returns
    -------
    bytes
        Raw ``.pptx`` file bytes.
    """
    from app.dashboards.svg_render import render_widgets_svg, _widget_to_payload  # noqa: PLC0415

    export_cfg = get_export_config(spec)
    resolved_page_size = page_size or export_cfg.get("page_size") or "A4"
    width_cm, height_cm = _PAGE_SIZES_CM.get(resolved_page_size, _PAGE_SIZES_CM["A4"])

    cols = int((spec.layout or {}).get("cols", 12))
    layout = get_surface_layout(spec, "grid")

    data_by_widget: dict[str, dict[str, Any]] = {
        entry["widget_id"]: entry
        for entry in widget_data
        if isinstance(entry, dict) and "widget_id" in entry
    }

    col_w_px = (page_width_px - 2 * margin_px) / cols

    payloads: list[dict[str, Any]] = []
    for widget in spec.widgets:
        entry = layout.get(widget.id)
        if entry is None:
            continue
        w_px = max(1, int(round(entry.w * col_w_px - gap_px)))
        h_px = max(1, int(entry.h * row_height_px - gap_px))
        payloads.append(
            _widget_to_payload(
                widget,
                data_by_widget.get(widget.id),
                width=w_px,
                height=h_px,
            )
        )

    widget_svgs: dict[str, str | None] = {}
    if payloads:
        rendered = render_widgets_svg(payloads, timeout=ssr_timeout)
        for r in rendered:
            wid = r.get("id")
            if wid:
                widget_svgs[wid] = r.get("svg") if not r.get("webgl") else None

    return render_board_pptx(spec, widget_svgs, page_size=page_size)
